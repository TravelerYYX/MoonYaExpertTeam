#!/usr/bin/env python3
"""Generate and reopen every Office artifact type required by File Agent."""

from pathlib import Path
from tempfile import TemporaryDirectory

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pypdf import PdfReader, PdfWriter


def main() -> None:
    with TemporaryDirectory(prefix="moonya-office-") as temp_dir:
        root = Path(temp_dir)

        docx_path = root / "verify.docx"
        document = Document()
        document.add_heading("MoonYa Office 验证", level=1)
        document.add_paragraph("docx content verified")
        document.save(docx_path)
        reopened_document = Document(docx_path)
        assert any(
            paragraph.text == "docx content verified"
            for paragraph in reopened_document.paragraphs
        )

        xlsx_path = root / "verify.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "验收"
        sheet["A1"] = "xlsx content verified"
        workbook.save(xlsx_path)
        reopened_workbook = load_workbook(xlsx_path, data_only=True)
        assert reopened_workbook.sheetnames == ["验收"]
        assert reopened_workbook["验收"]["A1"].value == "xlsx content verified"

        pptx_path = root / "verify.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "MoonYa Office 验证"
        slide.placeholders[1].text = "pptx content verified"
        presentation.save(pptx_path)
        reopened_presentation = Presentation(pptx_path)
        assert len(reopened_presentation.slides) == 1
        assert reopened_presentation.slides[0].shapes.title.text == "MoonYa Office 验证"

        pdf_path = root / "verify.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=595, height=842)
        writer.add_metadata({"/Title": "MoonYa PDF verification"})
        with pdf_path.open("wb") as output:
            writer.write(output)
        reopened_pdf = PdfReader(pdf_path)
        assert len(reopened_pdf.pages) == 1
        assert reopened_pdf.metadata.title == "MoonYa PDF verification"

        generated = [docx_path, xlsx_path, pptx_path, pdf_path]
        assert all(path.is_file() and path.stat().st_size > 0 for path in generated)
        print(
            "office runtime smoke: PASS "
            + ", ".join(f"{path.suffix[1:]}={path.stat().st_size}" for path in generated)
        )


if __name__ == "__main__":
    main()
